# data_ingestion/loader.py
"""
Purpose
-------
High-level loader orchestrator for rag-chatbot-docsearch.

Responsibilities (summary)
- Discover candidate files (paths, globs, upload file-likes).
- Produce canonical FileDescriptor dictionaries for discovered inputs.
- Choose appropriate Handler for each file type.
- Stream page-level raw extractions (text or images) as Document records for downstream processing.
- Provide utilities for temp-writing uploads, quick hashing, mime sniffing, and utf heuristics.

Notes
-----
This file intentionally mixes:
- small, safe helper implementations you can reuse immediately, and
- well-commented skeleton classes/methods you will implement step-by-step.

Design guidance:
- Keep extraction (loader/handlers) separate from cleaning/chunking.
- Cache expensive steps (rendering, OCR results) in your cache_dir.
- Persist per-file status/manifest in a small sqlite/jsonl store (StatusStore below is a placeholder).
"""
from pathlib import Path
from typing import Iterator, Dict, Optional, Any, Iterable, Union, List, Tuple
import tempfile
import hashlib
import mimetypes
import logging
import fnmatch
import os
import io
import time
import abc
import json

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

# Optional enhancement: if python-magic (libmagic) is installed we use it for better sniffing.
try:
    import magic  # type: ignore
    HAS_MAGIC = True
except Exception:
    HAS_MAGIC = False

# ---------------------------
# Configurable defaults
# ---------------------------
QUICK_HASH_WINDOW = 64 * 1024  # bytes read for quick hashing (64KB)
MIME_SNIFF_BYTES = 512
RECURSIVE_BY_DEFAULT = True
IGNORE_PATTERNS_DEFAULT: List[str] = ["*.DS_Store", "*.tmp", "*.~*"]

# extension groups
TEXT_EXTS = {".txt", ".md", ".csv", ".log"}
PDF_EXTS = {".pdf"}
DOCX_EXTS = {".docx", ".pptx", ".xlsx"}
ARCHIVE_EXTS = {".zip", ".tar", ".tar.gz", ".tgz"}
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".tiff", ".bmp"}

# ---------------------------
# Helper utilities (implemented)
# ---------------------------


def _is_file_like(obj: Any) -> bool:
    """
    Heuristic to decide whether obj is a readable file-like object.
    We prefer objects exposing read() (and optionally seek()).
    """
    return hasattr(obj, "read") and callable(getattr(obj, "read"))


def _read_head_bytes_from_path(path: Path, n: int) -> bytes:
    """Return first n bytes of a file at path. Return b'' on error."""
    try:
        with open(path, "rb") as f:
            return f.read(n)
    except Exception:
        return b""


def _read_head_bytes_from_filelike(fobj: io.IOBase, n: int) -> bytes:
    """
    Read first n bytes from a file-like. If seekable, restore position afterwards.
    This is safe for UploadFile / io.BytesIO objects.
    """
    try:
        pos = fobj.tell()
    except Exception:
        pos = None

    try:
        data = fobj.read(n) or b""
    except Exception:
        data = b""

    try:
        if pos is not None:
            fobj.seek(pos)
    except Exception:
        pass
    # If data is str, encode
    if isinstance(data, str):
        return data.encode("utf-8", errors="replace")
    return data


def _sniff_mime(bytes_head: bytes, path: Optional[Path] = None) -> str:
    """
    Best-effort mime-type sniffing.
    - If python-magic available, use it on bytes_head.
    - Otherwise fall back to mimetypes.guess_type(path).
    Returns empty string if unknown.
    """
    if HAS_MAGIC and bytes_head:
        try:
            m = magic.Magic(mime=True)
            return m.from_buffer(bytes_head) or ""
        except Exception:
            pass
    if path is not None:
        guess, _ = mimetypes.guess_type(str(path))
        return guess or ""
    return ""


def _quick_hash_of(path: Optional[Path], filelike: Optional[io.IOBase], window: int = QUICK_HASH_WINDOW) -> str:
    """
    Compute a fast 'quick hash' to detect changes or dedupe quickly.
    Strategy:
    - If path provided: read first `window` bytes, plus append file size and mtime for more entropy.
    - If filelike provided: read first `window` bytes from the stream (attempt to preserve position).
    Returns hex sha256 digest.
    """
    h = hashlib.sha256()
    if path is not None:
        try:
            with open(path, "rb") as f:
                head = f.read(window)
                h.update(head)
        except Exception:
            # still return a valid digest
            h.update(b"")
        try:
            stat = path.stat()
            h.update(str(stat.st_size).encode())
            h.update(str(int(stat.st_mtime)).encode())
        except Exception:
            pass
    elif filelike is not None:
        head = _read_head_bytes_from_filelike(filelike, window)
        h.update(head)
    else:
        h.update(b"")
    return h.hexdigest()


def _utf_likely_from_bytes(head: bytes) -> bool:
    """
    Heuristic whether bytes appear to be text (utf-8 friendly).
    - If many NUL bytes -> likely binary
    - Try decoding as utf-8; if success, True. If latin-1 decodes but utf fails, return False (text but not utf).
    """
    if not head:
        return True
    nul_count = head.count(b"\x00")
    if len(head) and (nul_count / len(head)) > 0.01:  # >1% NUL bytes
        return False
    try:
        head.decode("utf-8")
        return True
    except Exception:
        try:
            head.decode("latin-1")
            return False
        except Exception:
            return False


def _suggest_handler_from_ext_and_mime(ext: Optional[str], mime: Optional[str]) -> str:
    """
    Suggest a coarse handler key from extension and mime.
    Returns one of: 'txt', 'pdf', 'docx', 'archive', 'image', 'binary'
    """
    ext = (ext or "").lower()
    mime = (mime or "").lower()
    if ext in TEXT_EXTS:
        return "txt"
    if ext in PDF_EXTS or mime == "application/pdf":
        return "pdf"
    if ext in DOCX_EXTS:
        return "docx"
    if ext in ARCHIVE_EXTS or ("zip" in (mime or "")):
        return "archive"
    if ext in IMAGE_EXTS or (mime.startswith("image/")):
        return "image"
    if mime.startswith("text/"):
        return "txt"
    return "binary"


def _write_filelike_to_temp(filelike: io.IOBase, filename_hint: Optional[str], cache_dir: Optional[Union[str, Path]] = None) -> Path:
    """
    Persist a file-like object to a temporary file under cache_dir (or system tmp dir) and return the Path.
    - This simplifies integration with libraries that require file paths (PyMuPDF, pypdf, python-docx).
    - Caller should ensure cache_dir exists and handle cleanup policy.
    """
    parent = Path(cache_dir) if cache_dir is not None else Path(tempfile.gettempdir())
    parent.mkdir(parents=True, exist_ok=True)
    suffix = Path(filename_hint).suffix if filename_hint else ""
    with tempfile.NamedTemporaryFile(delete=False, dir=str(parent), suffix=suffix) as tmp:
        # try to reset to beginning when possible
        try:
            filelike.seek(0)
        except Exception:
            pass
        while True:
            chunk = filelike.read(64 * 1024)
            if not chunk:
                break
            if isinstance(chunk, str):
                chunk = chunk.encode("utf-8", errors="replace")
            tmp.write(chunk)
        tmp.flush()
        tmp_path = Path(tmp.name)
    return tmp_path


# ---------------------------
# Data contracts (lightweight)
# ---------------------------


def file_descriptor_template() -> Dict[str, Any]:
    """
    Returns the canonical FileDescriptor shape (helpful for tests and docs).
    """
    return {
        "source": None,  # original input: path string or upload identifier
        "path": None,  # Path on disk if available (temp path for uploads)
        "file_like": None,  # original file-like if not written to disk
        "ext": None,  # lower-case extension
        "mime": None,  # optional sniffer result
        "size": None,
        "mtime": None,
        "quick_hash": None,
        "suggested_handler": None,
        "notes": [],
    }


def document_template() -> Dict[str, Any]:
    """
    Document contract produced by loader/handlers (raw extraction stage).
    Document is minimal: text or image, id, and metadata for provenance.
    """
    return {
        "id": None,
        "raw_text": None,  # may be None if image-only and OCR not run here
        "image_bytes": None,  # optional; only for image pages or when rendering used
        "metadata": {
            "source_path": None,
            "file_type": None,
            "page_number": None,
            "file_hash": None,
            "page_hash": None,
            "extraction_engine": None,
            "ocr_engine": None,
            "ocr_confidence": None,
            "ingested_at": None,
            "offsets": None,
        },
    }


# ---------------------------
# Handler base classes & stubs
# ---------------------------


class Handler(abc.ABC):
    """
    Abstract handler base class.

    Concrete handlers should implement:
    - inspect(file_descriptor) -> dict (file-level metadata like page_count, has_embedded_text)
    - extract(file_descriptor) -> yields page-level extractions (page_number, raw_text, image_bytes, meta)

    Handler implementations should not perform cleaning; they produce raw extraction only.
    """

    @abc.abstractmethod
    def inspect(self, fd: Dict[str, Any]) -> Dict[str, Any]:
        """
        Quick inspection without heavy work.
        Example outputs:
          { "page_count": 3, "has_embedded_text": True, "languages": ["en"] }
        """
        raise NotImplementedError

    @abc.abstractmethod
    def extract(self, fd: Dict[str, Any]) -> Iterator[Dict[str, Any]]:
        """
        Yield page-level dicts:
          { "page_number": int, "raw_text": Optional[str], "image_bytes": Optional[bytes], "meta": {...} }
        Extraction must be streaming where possible to avoid reading entire file into memory.
        """
        raise NotImplementedError


# Concrete handler stubs: implement these gradually


class TxtHandler(Handler):
    """
    Handler for plain text files.
    Implementation notes:
    - Prefer opening in binary and decoding in chunks using detected encoding (charset-normalizer or chardet).
    - Stream lines to avoid memory spikes for large text files.
    - Yield a single 'page' or split into multi-line pages based on config (e.g., N lines per document).
    """

    def inspect(self, fd: Dict[str, Any]) -> Dict[str, Any]:
        # Minimal quick inspection for text: page_count = 1, has_embedded_text True
        return {"page_count": 1, "has_embedded_text": True}

    def extract(self, fd: Dict[str, Any]) -> Iterator[Dict[str, Any]]:
        """
        Example simple behaviour:
        - If fd['path'] exists: open and yield a single page with raw_text = file content (or chunked).
        - If fd['file_like'] provided and path not set: read from file_like (seek to 0 if possible).
        """
        path = fd.get("path")
        filelike = fd.get("file_like")
        head = b""
        if path:
            try:
                with open(path, "rb") as f:
                    head = f.read()
                    try:
                        raw = head.decode("utf-8")
                    except Exception:
                        raw = head.decode("latin-1", errors="replace")
                    yield {"page_number": 1, "raw_text": raw, "image_bytes": None, "meta": {}}
            except Exception as e:
                logger.exception("TxtHandler failed to read path %s: %s", path, e)
                return
        elif filelike:
            try:
                try:
                    filelike.seek(0)
                except Exception:
                    pass
                data = filelike.read()
                if isinstance(data, bytes):
                    try:
                        raw = data.decode("utf-8")
                    except Exception:
                        raw = data.decode("latin-1", errors="replace")
                else:
                    raw = str(data)
                yield {"page_number": 1, "raw_text": raw, "image_bytes": None, "meta": {}}
            except Exception as e:
                logger.exception("TxtHandler failed to read filelike: %s", e)
                return
        else:
            return


class PdfHandler(Handler):
    """
    PDF handler (text-first strategy).
    Implementation plan:
    - inspect(): try to use a light-weight PDF library to return page_count and whether pages have embedded_text
    - extract(): iterate pages:
        * If embedded text present for page -> yield raw_text
        * Else -> yield image_bytes (rendered page) with meta indicating OCR fallback needed
    Notes:
    - Many PDF libraries require a filesystem path; if so, ensure discover wrote uploads to temp files.
    - Keep heavy rendering and OCR out of extract() unless configured to run OCR.
    """

    def inspect(self, fd: Dict[str, Any]) -> Dict[str, Any]:
        """Inspect PDF file to get page count and text availability"""
        try:
            import pypdf
            with open(fd['path'], 'rb') as file:
                reader = pypdf.PdfReader(file)
                page_count = len(reader.pages)
                
                # Check if pages have embedded text
                has_embedded_text = any(
                    len(page.extract_text().strip()) > 0 
                    for page in reader.pages[:min(3, page_count)]  # Check first 3 pages
                )
                
                return {
                    "page_count": page_count,
                    "has_embedded_text": has_embedded_text
                }
        except Exception as e:
            logger.error(f"Error inspecting PDF {fd['path']}: {e}")
            return {"page_count": None, "has_embedded_text": None}

    def extract(self, fd: Dict[str, Any]) -> Iterator[Dict[str, Any]]:
        """Extract text from PDF pages"""
        try:
            import pypdf
            with open(fd['path'], 'rb') as file:
                reader = pypdf.PdfReader(file)
                
                for page_num, page in enumerate(reader.pages):
                    try:
                        text = page.extract_text()
                        if text.strip():
                            yield {
                                'content': text,
                                'page_number': page_num + 1,
                                'source': str(fd['path']),
                                'content_type': 'text',
                                'metadata': {
                                    'file_path': str(fd['path']),
                                    'page_number': page_num + 1,
                                    'total_pages': len(reader.pages)
                                }
                            }
                    except Exception as e:
                        logger.warning(f"Error extracting text from page {page_num + 1}: {e}")
                        continue
        except Exception as e:
            logger.error(f"Error extracting PDF {fd['path']}: {e}")
            return


class DocxHandler(Handler):
    """
    Office handler for DOCX/PPTX/XLSX where appropriate.
    Notes:
    - Use python-docx for docx paragraphs
    - For pptx, use python-pptx to extract text per slide
    - Yield sensible page_number / paragraph indexing for provenance
    """

    def inspect(self, fd: Dict[str, Any]) -> Dict[str, Any]:
        """Inspect DOCX/PPTX file to get structure info"""
        try:
            ext = fd.get('ext', '').lower()
            if ext == '.docx':
                from docx import Document
                doc = Document(fd['path'])
                paragraph_count = len(doc.paragraphs)
                return {
                    "page_count": paragraph_count,  # Approximate
                    "has_embedded_text": paragraph_count > 0
                }
            elif ext == '.pptx':
                from pptx import Presentation
                prs = Presentation(fd['path'])
                slide_count = len(prs.slides)
                return {
                    "page_count": slide_count,
                    "has_embedded_text": slide_count > 0
                }
            else:
                return {"page_count": None, "has_embedded_text": None}
        except Exception as e:
            logger.error(f"Error inspecting office file {fd['path']}: {e}")
            return {"page_count": None, "has_embedded_text": None}

    def extract(self, fd: Dict[str, Any]) -> Iterator[Dict[str, Any]]:
        """Extract text from DOCX/PPTX files"""
        try:
            ext = fd.get('ext', '').lower()
            
            if ext == '.docx':
                from docx import Document
                doc = Document(fd['path'])
                
                for para_num, paragraph in enumerate(doc.paragraphs):
                    text = paragraph.text.strip()
                    if text:
                        yield {
                            'content': text,
                            'page_number': para_num + 1,
                            'source': str(fd['path']),
                            'content_type': 'text',
                            'metadata': {
                                'file_path': str(fd['path']),
                                'paragraph_number': para_num + 1,
                                'total_paragraphs': len(doc.paragraphs)
                            }
                        }
            
            elif ext == '.pptx':
                from pptx import Presentation
                prs = Presentation(fd['path'])
                
                for slide_num, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    if slide_text:
                        yield {
                            'content': '\n'.join(slide_text),
                            'page_number': slide_num + 1,
                            'source': str(fd['path']),
                            'content_type': 'text',
                            'metadata': {
                                'file_path': str(fd['path']),
                                'slide_number': slide_num + 1,
                                'total_slides': len(prs.slides)
                            }
                        }
            
        except Exception as e:
            logger.error(f"Error extracting office file {fd['path']}: {e}")
            return


class ImageHandler(Handler):
    """
    Image handler: returns image_bytes for OCR downstream.
    - Do not run OCR here unless explicitly configured. Just yield the bytes and minimal metadata.
    """

    def inspect(self, fd: Dict[str, Any]) -> Dict[str, Any]:
        return {"page_count": 1, "has_embedded_text": False}

    def extract(self, fd: Dict[str, Any]) -> Iterator[Dict[str, Any]]:
        path = fd.get("path")
        filelike = fd.get("file_like")
        if path:
            try:
                b = path.read_bytes()
                yield {"page_number": 1, "raw_text": None, "image_bytes": b, "meta": {}}
            except Exception as e:
                logger.exception("ImageHandler failed reading path %s: %s", path, e)
                return
        elif filelike:
            try:
                try:
                    filelike.seek(0)
                except Exception:
                    pass
                data = filelike.read()
                if isinstance(data, str):
                    data = data.encode("utf-8", errors="replace")
                yield {"page_number": 1, "raw_text": None, "image_bytes": data, "meta": {}}
            except Exception as e:
                logger.exception("ImageHandler failed reading filelike: %s", e)
                return
        else:
            return


class ArchiveHandler(Handler):
    """
    Archive handler: iterate entries inside a zip/tar and delegate each entry to the appropriate handler.
    Important: track compound source_path like "archive.zip::folder/file.pdf" for provenance.
    """

    def inspect(self, fd: Dict[str, Any]) -> Dict[str, Any]:
        """Inspect archive to get entry count"""
        try:
            import zipfile
            import tarfile
            
            path = fd['path']
            entry_count = 0
            
            if path.suffix.lower() == '.zip':
                with zipfile.ZipFile(path, 'r') as zf:
                    entry_count = len(zf.namelist())
            elif path.suffix.lower() in ['.tar', '.gz'] or path.name.endswith('.tar.gz'):
                with tarfile.open(path, 'r') as tf:
                    entry_count = len(tf.getnames())
            
            return {
                "page_count": entry_count,
                "has_embedded_text": entry_count > 0
            }
        except Exception as e:
            logger.error(f"Error inspecting archive {fd['path']}: {e}")
            return {"page_count": None, "has_embedded_text": None}

    def extract(self, fd: Dict[str, Any]) -> Iterator[Dict[str, Any]]:
        """Extract and process archive entries"""
        try:
            import zipfile
            import tarfile
            import tempfile
            from pathlib import Path
            
            path = fd['path']
            temp_dir = tempfile.mkdtemp(prefix="archive_extract_")
            
            try:
                if path.suffix.lower() == '.zip':
                    with zipfile.ZipFile(path, 'r') as zf:
                        for entry_name in zf.namelist():
                            if not entry_name.endswith('/'):  # Skip directories
                                # Extract to temp file
                                temp_path = Path(temp_dir) / entry_name
                                temp_path.parent.mkdir(parents=True, exist_ok=True)
                                
                                with zf.open(entry_name) as source, open(temp_path, 'wb') as target:
                                    target.write(source.read())
                                
                                # Create file descriptor for extracted file
                                extracted_fd = {
                                    'source': f"{path}::{entry_name}",
                                    'path': temp_path,
                                    'ext': Path(entry_name).suffix.lower(),
                                    'size': temp_path.stat().st_size,
                                    'mtime': temp_path.stat().st_mtime,
                                    'notes': ['extracted_from_archive']
                                }
                                
                                # Delegate to appropriate handler
                                handler = self._get_handler_for_file(extracted_fd)
                                if handler:
                                    for result in handler.extract(extracted_fd):
                                        # Update source path to include archive info
                                        result['source'] = f"{path}::{entry_name}"
                                        result['metadata'] = result.get('metadata', {})
                                        result['metadata']['archive_path'] = str(path)
                                        result['metadata']['archive_entry'] = entry_name
                                        yield result
                
                elif path.suffix.lower() in ['.tar', '.gz'] or path.name.endswith('.tar.gz'):
                    with tarfile.open(path, 'r') as tf:
                        for member in tf.getmembers():
                            if member.isfile():
                                # Extract to temp file
                                temp_path = Path(temp_dir) / member.name
                                temp_path.parent.mkdir(parents=True, exist_ok=True)
                                
                                with tf.extractfile(member) as source, open(temp_path, 'wb') as target:
                                    target.write(source.read())
                                
                                # Create file descriptor for extracted file
                                extracted_fd = {
                                    'source': f"{path}::{member.name}",
                                    'path': temp_path,
                                    'ext': Path(member.name).suffix.lower(),
                                    'size': member.size,
                                    'mtime': member.mtime,
                                    'notes': ['extracted_from_archive']
                                }
                                
                                # Delegate to appropriate handler
                                handler = self._get_handler_for_file(extracted_fd)
                                if handler:
                                    for result in handler.extract(extracted_fd):
                                        # Update source path to include archive info
                                        result['source'] = f"{path}::{member.name}"
                                        result['metadata'] = result.get('metadata', {})
                                        result['metadata']['archive_path'] = str(path)
                                        result['metadata']['archive_entry'] = member.name
                                        yield result
            
            finally:
                # Clean up temp directory
                import shutil
                shutil.rmtree(temp_dir, ignore_errors=True)
                
        except Exception as e:
            logger.error(f"Error extracting archive {fd['path']}: {e}")
            return
    
    def _get_handler_for_file(self, fd: Dict[str, Any]) -> Optional[Handler]:
        """Get appropriate handler for extracted file"""
        ext = fd.get('ext', '').lower()
        
        if ext in TEXT_EXTS:
            return TxtHandler()
        elif ext in PDF_EXTS:
            return PdfHandler()
        elif ext in DOCX_EXTS:
            return DocxHandler()
        elif ext in IMAGE_EXTS:
            return ImageHandler()
        else:
            logger.warning(f"No handler for file type: {ext}")
            return None


# ---------------------------
# OCR client and cache stub
# ---------------------------


class OcrClient:
    """
    Minimal OCR client interface.
    Implementations should provide .process(image_bytes, lang_hint=None) returning:
      { "text": str, "boxes": optional list, "confidence": float, "engine": "tesseract" or "gcloud", "meta": {...} }

    Concrete implementations should also expose engine version for provenance.
    """

    def __init__(self, provider: str = "local", config: Optional[Dict[str, Any]] = None):
        self.provider = provider
        self.config = config or {}

    def process(self, image_bytes: bytes, lang_hint: Optional[str] = None) -> Dict[str, Any]:
        """Process image bytes with OCR"""
        try:
            if self.provider == "local":
                return self._process_local(image_bytes, lang_hint)
            elif self.provider == "cloud":
                return self._process_cloud(image_bytes, lang_hint)
            else:
                raise ValueError(f"Unknown OCR provider: {self.provider}")
        except Exception as e:
            logger.error(f"Error processing OCR: {e}")
            return {
                "text": "",
                "boxes": [],
                "confidence": 0.0,
                "engine": "error",
                "meta": {"error": str(e)}
            }
    
    def _process_local(self, image_bytes: bytes, lang_hint: Optional[str] = None) -> Dict[str, Any]:
        """Process with local Tesseract OCR"""
        try:
            import pytesseract
            from PIL import Image
            import io
            
            # Convert bytes to PIL Image
            image = Image.open(io.BytesIO(image_bytes))
            
            # Configure Tesseract
            config = '--oem 3 --psm 6'
            if lang_hint:
                config += f' -l {lang_hint}'
            
            # Extract text
            text = pytesseract.image_to_string(image, config=config)
            
            # Get confidence data
            try:
                data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                confidences = [int(conf) for conf in data['conf'] if int(conf) > 0]
                avg_confidence = sum(confidences) / len(confidences) if confidences else 0
            except:
                avg_confidence = 0.0
            
            return {
                "text": text.strip(),
                "boxes": [],  # Could be implemented with pytesseract.image_to_boxes
                "confidence": avg_confidence / 100.0,  # Normalize to 0-1
                "engine": "tesseract",
                "meta": {
                    "provider": "local",
                    "language_hint": lang_hint,
                    "image_size": image.size
                }
            }
        except ImportError:
            logger.error("pytesseract not installed. Install with: pip install pytesseract")
            return {
                "text": "",
                "boxes": [],
                "confidence": 0.0,
                "engine": "tesseract",
                "meta": {"error": "pytesseract not installed"}
            }
    
    def _process_cloud(self, image_bytes: bytes, lang_hint: Optional[str] = None) -> Dict[str, Any]:
        """Process with cloud OCR service (placeholder)"""
        # This would integrate with Google Cloud Vision, AWS Textract, etc.
        logger.warning("Cloud OCR not implemented yet")
        return {
            "text": "",
            "boxes": [],
            "confidence": 0.0,
            "engine": "cloud",
            "meta": {"error": "Cloud OCR not implemented"}
        }


# ---------------------------
# Status/manifest store stub
# ---------------------------


class StatusStore:
    """
    Simple status store interface for per-file processing metadata and idempotency.
    Implementations could be:
      - jsonl files under cache_dir/manifests
      - sqlite DB (recommended for scale)
      - Redis (if distributed)
    Minimal required methods:
      - get_status(file_hash) -> dict or None
      - mark_processed(file_hash, details)
      - mark_failed(file_hash, details)
    """

    def __init__(self, cache_dir: Optional[Union[str, Path]] = None):
        self.cache_dir = Path(cache_dir) if cache_dir else None
        # TODO: implement persistent backing store
        self._store: Dict[str, Dict[str, Any]] = {}

    def get_status(self, file_hash: str) -> Optional[Dict[str, Any]]:
        return self._store.get(file_hash)

    def mark_processed(self, file_hash: str, details: Dict[str, Any]):
        self._store[file_hash] = {"state": "processed", "details": details, "timestamp": int(time.time())}

    def mark_failed(self, file_hash: str, details: Dict[str, Any]):
        self._store[file_hash] = {"state": "failed", "details": details, "timestamp": int(time.time())}


# ---------------------------
# Loader orchestrator
# ---------------------------


class Loader:
    """
    Loader orchestrates discovery and extraction.

    Responsibilities in methods below:
    - discover: implemented as a solid, usable function (it reuses helper utilities).
    - choose_handler: pick a Handler instance for a FileDescriptor.
    - process_file: orchestrate per-file work, yield Document dicts.
    - stream_documents: top-level convenience to process many inputs.
    - preview: return first N document samples for a file (dry-run friendly).
    """

    def __init__(
        self,
        *,
        cache_dir: Optional[Union[str, Path]] = None,
        handlers: Optional[Dict[str, Handler]] = None,
        ocr_client: Optional[OcrClient] = None,
        status_store: Optional[StatusStore] = None,
        ignore_patterns: Optional[Iterable[str]] = None,
        recursive: bool = RECURSIVE_BY_DEFAULT,
        # optional storage adapters (additive; leave None to keep behaviour unchanged)
        blob_store: Optional[Any] = None,        # expected API: put_blob(bucket, key, bytes) -> path/uri
        doc_store: Optional[Any] = None,         # expected API: upsert_document(doc_dict)
        vector_client: Optional[Any] = None,     # optional vector adapter (ChromaAdapter)
        index_on_ingest: bool = False,           # whether to attempt indexing embeddings while ingesting
    ):
        """
        Setup loader state and handler registry.

        Arguments:
        - cache_dir: local dir for temp files, OCR cache, manifests
        - handlers: mapping of handler_key -> Handler instance (e.g., "txt": TxtHandler())
        - ocr_client: optional OCR client for fallback
        - status_store: store for processed file metadata
        - ignore_patterns: list of glob patterns to skip
        - recursive: whether to walk directories recursively
        - blob_store: optional filesystem/S3-like blob store adapter
        - doc_store: optional metadata store (sqlite wrapper) exposing upsert_document/get_document_by_id
        - vector_client: optional vector adapter (e.g., ChromaAdapter)
        - index_on_ingest: if True, attempts to index embeddings present in metadata on ingest
        """
        self.cache_dir = Path(cache_dir) if cache_dir is not None else None
        self.handlers = handlers or {}
        # ensure default handlers are present
        self.handlers.setdefault("txt", TxtHandler())
        self.handlers.setdefault("pdf", PdfHandler())
        self.handlers.setdefault("docx", DocxHandler())
        self.handlers.setdefault("image", ImageHandler())
        self.handlers.setdefault("archive", ArchiveHandler())
        self.ocr_client = ocr_client
        self.status_store = status_store or StatusStore(cache_dir=self.cache_dir)
        self.ignore_patterns = list(ignore_patterns) if ignore_patterns else IGNORE_PATTERNS_DEFAULT
        self.recursive = recursive

        # optional storage wiring
        self.blob_store = blob_store
        self.doc_store = doc_store
        self.vector_client = vector_client
        self.index_on_ingest = index_on_ingest

    # -----------------------
    # discover (implemented)
    # -----------------------
    def discover(
        self,
        input_path_or_filelike: Union[str, Path, Iterable[Union[str, Path, Any]], Any],
        *,
        recursive: Optional[bool] = None,
        ignore_patterns: Optional[Iterable[str]] = None,
        quick_hash_window: int = QUICK_HASH_WINDOW,
    ) -> Iterator[Dict[str, Any]]:
        """
        Discover files and yield canonical FileDescriptor dictionaries.

        This method handles:
        - single paths, lists of paths, glob patterns, directories
        - file-like upload objects (writes to a temp path in cache_dir)
        - ignore patterns and recursive directory walking

        Best practice: call this with a path, then pass each descriptor to process_file().
        """
        recursive = self.recursive if recursive is None else recursive
        ignore_patterns = list(ignore_patterns) if ignore_patterns is not None else self.ignore_patterns

        items: Iterable[Any]
        if input_path_or_filelike is None:
            return
        if isinstance(input_path_or_filelike, (str, Path)) or _is_file_like(input_path_or_filelike):
            items = [input_path_or_filelike]
        elif isinstance(input_path_or_filelike, Iterable):
            items = input_path_or_filelike  # type: ignore
        else:
            items = [input_path_or_filelike]

        for item in items:
            # Handle glob strings quickly
            try:
                if isinstance(item, str) and any(ch in item for ch in ["*", "?", "["]):
                    for p in Path().glob(item):
                        if any(fnmatch.fnmatch(str(p.name), pat) for pat in ignore_patterns):
                            continue
                        if p.is_dir():
                            for sub in p.rglob("*") if recursive else p.glob("*"):
                                if sub.is_file():
                                    yield from self.discover(sub, recursive=recursive, ignore_patterns=ignore_patterns, quick_hash_window=quick_hash_window)
                            continue
                        yield from self.discover(p, recursive=recursive, ignore_patterns=ignore_patterns, quick_hash_window=quick_hash_window)
                    continue
            except Exception as e:
                logger.exception("Error expanding glob %s: %s", item, e)

            # Path-like handling
            if isinstance(item, (str, Path)) and not _is_file_like(item):
                path = Path(item)
                if path.is_dir():
                    for sub in path.rglob("*") if recursive else path.glob("*"):
                        if not sub.is_file():
                            continue
                        if any(fnmatch.fnmatch(str(sub.name), pat) for pat in ignore_patterns):
                            continue
                        yield from self.discover(sub, recursive=recursive, ignore_patterns=ignore_patterns, quick_hash_window=quick_hash_window)
                    continue

                fd = file_descriptor_template()
                fd["source"] = str(item)
                fd["path"] = path
                fd["ext"] = path.suffix.lower()
                fd["notes"] = []
                try:
                    try:
                        stat = path.stat()
                        fd["size"] = int(stat.st_size)
                        fd["mtime"] = int(stat.st_mtime)
                    except Exception:
                        fd["notes"].append("stat_failed")

                    head = _read_head_bytes_from_path(path, MIME_SNIFF_BYTES)
                    fd["mime"] = _sniff_mime(head, path=path)
                    fd["quick_hash"] = _quick_hash_of(path=path, filelike=None, window=quick_hash_window)
                    fd["notes"].append("utf_likely" if _utf_likely_from_bytes(head) else "non_utf_likely")
                    fd["suggested_handler"] = _suggest_handler_from_ext_and_mime(fd["ext"], fd.get("mime", ""))
                except Exception as e:
                    logger.exception("Error inspecting path %s: %s", item, e)
                    fd["notes"].append(f"inspect_error:{e!s}")
                yield fd
                continue

            # File-like (upload) handling
            if _is_file_like(item):
                filelike = item
                filename_hint = getattr(filelike, "filename", None) or getattr(filelike, "name", None) or "upload"
                fd = file_descriptor_template()
                fd["source"] = getattr(filelike, "filename", None) or "<file-like>"
                fd["file_like"] = filelike
                fd["ext"] = Path(filename_hint).suffix.lower()
                fd["notes"].append("uploaded")
                try:
                    head = _read_head_bytes_from_filelike(filelike, MIME_SNIFF_BYTES)
                    fd["mime"] = _sniff_mime(head, path=Path(filename_hint) if filename_hint else None)
                    # try to guess size if seekable
                    try:
                        filelike.seek(0, os.SEEK_END)
                        fd["size"] = filelike.tell()
                        filelike.seek(0)
                    except Exception:
                        fd["size"] = None
                    fd["quick_hash"] = _quick_hash_of(path=None, filelike=filelike, window=quick_hash_window)
                    fd["notes"].append("utf_likely" if _utf_likely_from_bytes(head) else "non_utf_likely")
                    fd["suggested_handler"] = _suggest_handler_from_ext_and_mime(fd["ext"], fd.get("mime", ""))
                    # write to temp path to simplify downstream libraries that require a path
                    try:
                        tmp_path = _write_filelike_to_temp(filelike, filename_hint, cache_dir=self.cache_dir)
                        fd["path"] = tmp_path
                        fd["notes"].append("temp_written")
                        try:
                            stat = tmp_path.stat()
                            fd["size"] = int(stat.st_size)
                            fd["mtime"] = int(stat.st_mtime)
                        except Exception:
                            pass
                    except Exception as e:
                        logger.exception("Failed to write upload to temp file: %s", e)
                        fd["notes"].append(f"temp_write_failed:{e!s}")
                except Exception as e:
                    logger.exception("Error inspecting upload: %s", e)
                    fd["notes"].append(f"upload_inspect_error:{e!s}")
                yield fd
                continue

            # Fallback unknown type
            try:
                yield {
                    "source": repr(item),
                    "path": None,
                    "file_like": None,
                    "ext": None,
                    "mime": None,
                    "size": None,
                    "mtime": None,
                    "quick_hash": None,
                    "suggested_handler": None,
                    "notes": ["unhandled_input_type"],
                }
            except Exception:
                continue

    # end discover

    # -----------------------
    # remaining orchestrator methods (implemented)
    # -----------------------

    def choose_handler(self, fd: Dict[str, Any]) -> Handler:
        """
        Decide which Handler instance to use for this FileDescriptor.

        Strategy:
        1. Use fd['suggested_handler'] if present and self.handlers contains it.
        2. Otherwise, use mime-based rules to pick a handler.
        3. Fallback to a generic 'binary' or 'txt' handler depending on utf heuristic.
        4. Raise or return a default handler if none match.

        Notes:
        - Keep this method deterministic (no heavy I/O).
        - Do not run OCR or rendering here.
        """
        suggested = fd.get("suggested_handler")
        if suggested and suggested in self.handlers:
            return self.handlers[suggested]
        mime = (fd.get("mime") or "").lower()
        ext = (fd.get("ext") or "").lower()
        # coarse rules
        if ext in TEXT_EXTS or mime.startswith("text"):
            return self.handlers["txt"]
        if ext in PDF_EXTS or "pdf" in mime:
            return self.handlers["pdf"]
        if ext in DOCX_EXTS:
            return self.handlers["docx"]
        if ext in IMAGE_EXTS or (mime.startswith("image")):
            return self.handlers["image"]
        if ext in ARCHIVE_EXTS or "zip" in mime:
            return self.handlers["archive"]
        # fallback
        return self.handlers.get("txt", list(self.handlers.values())[0])

    def _make_document_id(self, fd: Dict[str, Any], page_number: int, chunk_index: Optional[int] = None) -> str:
        """
        Helper for stable document ids. Example scheme:
          {file_quick_hash}::p{page_number}::c{chunk_index}
        Ensure id is deterministic based on file quick_hash and page index.
        """
        file_hash = fd.get("quick_hash") or "nofhash"
        if chunk_index is not None:
            return f"{file_hash}::p{page_number}::c{chunk_index}"
        return f"{file_hash}::p{page_number}"

    def _page_hash(self, page_number: int, raw_text: Optional[str], image_bytes: Optional[bytes]) -> str:
        """
        Compute a stable hash for a page using text if available, otherwise image bytes.
        """
        h = hashlib.sha256()
        if raw_text is not None:
            h.update(raw_text.encode("utf-8", errors="replace"))
        elif image_bytes is not None:
            h.update(image_bytes)
        h.update(str(page_number).encode())
        return h.hexdigest()

    def process_file(self, fd: Dict[str, Any], *, preview: bool = False, ocr_enabled: bool = True, force: bool = False, preview_limit: int = 2) -> Iterator[Dict[str, Any]]:
        """
        Orchestrate extraction for a single FileDescriptor and yield Document dicts.

        Behaviour summary:
        - Check status_store for existing processed state: skip unless force.
        - Choose handler and call inspect(fd).
        - Iterate handler.extract(fd):
            - compute page_hash
            - create Document dict using document_template
            - if page has image_bytes and ocr_enabled and ocr_client present: call OCR and attach
        - Update status_store.mark_processed on success or mark_failed on exceptions.
        - If preview True: stop after preview_limit pages and do not mark processed.
        """
        file_hash = fd.get("quick_hash")
        if not file_hash:
            file_hash = hashlib.sha256(str(fd.get("source")).encode()).hexdigest()

        # quick idempotency check
        try:
            status = self.status_store.get_status(file_hash) if self.status_store else None
            if status and status.get("state") == "processed" and not force and not preview:
                logger.info("Skipping already processed file %s (hash=%s)", fd.get("source"), file_hash)
                return
        except Exception:
            logger.exception("Error checking status store for %s", fd.get("source"))

        handler = None
        start_ts = time.time()
        pages_emitted = 0
        bytes_processed = 0
        try:
            handler = self.choose_handler(fd)
            logger.debug("Chosen handler %s for %s", handler.__class__.__name__, fd.get("source"))
            try:
                inspect_meta = handler.inspect(fd)
            except Exception as e:
                inspect_meta = {}
                logger.exception("Handler.inspect failed for %s: %s", fd.get("source"), e)

            extraction_engine = inspect_meta.get("engine") or handler.__class__.__name__

            for page_dict in handler.extract(fd):
                # guard: page_dict shape
                page_number = page_dict.get("page_number", 1)
                raw_text = page_dict.get("raw_text")
                image_bytes = page_dict.get("image_bytes")
                meta = page_dict.get("meta", {}) or {}

                page_hash = self._page_hash(page_number, raw_text, image_bytes)
                doc = document_template()
                doc["id"] = self._make_document_id(fd, page_number)
                doc["raw_text"] = raw_text
                doc["image_bytes"] = image_bytes
                metadata = doc["metadata"]
                metadata.update({
                    "source_path": str(fd.get("path") or fd.get("source")),
                    "file_type": fd.get("ext") or fd.get("mime"),
                    "page_number": page_number,
                    "file_hash": file_hash,
                    "page_hash": page_hash,
                    "extraction_engine": extraction_engine,
                    "ingested_at": int(time.time()),
                })
                # attach page-level meta
                metadata.update(meta)

                # If image and OCR desired, attempt OCR
                if image_bytes is not None and ocr_enabled and self.ocr_client is not None:
                    try:
                        ocr_result = self.ocr_client.process(image_bytes, lang_hint=meta.get("lang"))
                        if isinstance(ocr_result, dict):
                            ocr_text = ocr_result.get("text")
                            doc["raw_text"] = ocr_text or doc["raw_text"]
                            metadata["ocr_engine"] = ocr_result.get("engine")
                            metadata["ocr_confidence"] = ocr_result.get("confidence")
                            metadata["ocr_meta"] = ocr_result.get("meta")
                    except Exception:
                        logger.exception("OCR failed for %s page %s", fd.get("source"), page_number)

                # compute some counters
                if raw_text is not None:
                    bytes_processed += len(raw_text.encode("utf-8", errors="replace"))
                elif image_bytes is not None:
                    bytes_processed += len(image_bytes)

                # ------------- persist to storage (if provided) -------------
                # Only persist when not previewing (preview should be read-only / dry-run).
                if not preview:
                    # Persist image bytes (blobs)
                    if image_bytes is not None and getattr(self, "blob_store", None) is not None:
                        try:
                            # deterministic key: doc id plus extension
                            blob_key = f"{doc['id']}.bin"
                            # use file_hash as the bucket to group resources for a file
                            blob_path = self.blob_store.put_blob(str(file_hash), blob_key, image_bytes)
                            metadata["blob_path"] = blob_path
                            metadata["blob_key"] = blob_key
                        except Exception:
                            logger.exception("Failed to persist blob for %s page %s", fd.get("source"), page_number)

                    # Persist metadata document
                    if getattr(self, "doc_store", None) is not None:
                        try:
                            # Upsert the full document (IDs, raw_text, metadata). Avoid storing large image bytes in sqlite.
                            doc_to_store = dict(doc)
                            # keep raw_text and metadata; strip image bytes to avoid DB bloat
                            doc_to_store.pop("image_bytes", None)
                            self.doc_store.upsert_document(doc_to_store)
                        except Exception:
                            logger.exception("Failed to persist document metadata for %s page %s", fd.get("source"), page_number)

                    # Optional: index into vector DB (embedding must be provided by separate step)
                    if getattr(self, "vector_client", None) is not None and getattr(self, "index_on_ingest", False):
                        try:
                            emb = metadata.get("embedding")
                            if emb:
                                # vector_client adapter should expose add_documents(ids, embeddings, metadatas, documents)
                                self.vector_client.add_documents(ids=[doc["id"]], embeddings=[emb], metadatas=[metadata], documents=[doc.get("raw_text")])
                        except Exception:
                            logger.exception("Failed to index document %s on ingest", doc["id"])

                pages_emitted += 1
                yield doc

                # preview behaviour: early stop and do not mark processed
                if preview and pages_emitted >= preview_limit:
                    logger.debug("Preview stop after %s pages for %s", pages_emitted, fd.get("source"))
                    return

        except Exception as e:
            logger.exception("Processing failed for %s: %s", fd.get("source"), e)
            # mark failed if not preview
            if not preview:
                try:
                    self.status_store.mark_failed(file_hash, {"error": str(e), "source": fd.get("source")})
                except Exception:
                    logger.exception("Failed to mark_failed for %s", fd.get("source"))
            return
        finally:
            runtime = time.time() - start_ts
            if not preview:
                # record processed summary
                try:
                    summary = {"pages": pages_emitted, "bytes": bytes_processed, "runtime_seconds": runtime, "source": fd.get("source")}
                    self.status_store.mark_processed(file_hash, summary)
                except Exception:
                    logger.exception("Failed to mark_processed for %s", fd.get("source"))

    def stream_documents(self, inputs: Union[str, Path, Iterable[Any], Any], *, preview: bool = False, ocr_enabled: bool = True, force: bool = False) -> Iterator[Dict[str, Any]]:
        """
        Convenience top-level method:
        - Calls discover(inputs) to get FileDescriptors
        - For each descriptor, call process_file(fd) and yield all Documents
        - Collect and log basic metrics (files processed, pages extracted)

        Keep this method thin so you can call it from CLI or tests.
        """
        files_seen = 0
        total_pages = 0
        total_bytes = 0
        for fd in self.discover(inputs):
            files_seen += 1
            logger.info("Processing discovered file %s", fd.get("source"))
            try:
                for doc in self.process_file(fd, preview=preview, ocr_enabled=ocr_enabled, force=force):
                    total_pages += 1
                    # rough bytes: use metadata if present
                    try:
                        total_bytes += len(doc.get("raw_text", "") or b"") if isinstance(doc.get("raw_text"), (str, bytes)) else 0
                    except Exception:
                        pass
                    yield doc
            except Exception:
                logger.exception("Error while streaming documents for %s", fd.get("source"))
                continue
        logger.info("stream_documents summary: files=%s pages=%s bytes_approx=%s", files_seen, total_pages, total_bytes)

    def preview(self, input_item: Union[str, Path, Any], n_pages: int = 2) -> List[Dict[str, Any]]:
        """
        Extract first n_pages from the first file discovered for quick debugging.

        Behaviour:
        - Call discover(input_item) and pick the first descriptor
        - Call process_file(fd, preview=True) and collect up to n_pages documents
        - Return list of document dicts (do not persist status_store changes in preview mode)
        """
        docs = []
        it = self.discover(input_item)
        first_fd = None
        for fd in it:
            first_fd = fd
            break
        if first_fd is None:
            return docs
        try:
            for doc in self.process_file(first_fd, preview=True, ocr_enabled=False, force=True, preview_limit=n_pages):
                docs.append(doc)
                if len(docs) >= n_pages:
                    break
        except Exception:
            logger.exception("Preview failed for %s", first_fd.get("source"))
        return docs

    def mark_processed(self, fd: Dict[str, Any], details: Dict[str, Any]):
        """
        Convenience wrapper to mark file processed in status_store.
        details should contain summary info such as pages_count, bytes_processed, pipeline_tag.
        """
        file_hash = fd.get("quick_hash")
        if not file_hash:
            # fall back to path-based hash if needed
            file_hash = hashlib.sha256(str(fd.get("source")).encode()).hexdigest()
        self.status_store.mark_processed(file_hash, details)

    def get_status(self, fd: Dict[str, Any]) -> Optional[Dict[str, Any]]:
        """
        Get stored processing status for a file descriptor (if any).
        """
        file_hash = fd.get("quick_hash")
        if not file_hash:
            file_hash = hashlib.sha256(str(fd.get("source")).encode()).hexdigest()
        return self.status_store.get_status(file_hash)


# End of loader.py
